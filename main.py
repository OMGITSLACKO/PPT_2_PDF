import os
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import fitz

def ppt_to_pdf(ppt_path, pdf_path):
    prs = Presentation(ppt_path)
    c = canvas.Canvas(pdf_path, pagesize=letter)
    width, height = letter

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
                c.drawString(100, height - 100, text)
        c.showPage()
    c.save()

def merge_pdfs(pdf_list, output_path):
    if not pdf_list:
        print("No PDF files to merge.")
        return

    merged_pdf = fitz.open()
    for pdf in pdf_list:
        with fitz.open(pdf) as mfile:
            merged_pdf.insert_pdf(mfile)
    merged_pdf.save(output_path)

def list_files(directory):
    ppt_files = [f for f in os.listdir(directory) if f.endswith('.pptx')]
    pdf_files = [f for f in os.listdir(directory) if f.endswith('.pdf')]
    return ppt_files, pdf_files

def main():
    directory = input("Enter the directory containing PPT and PDF files: ")
    ppt_files, pdf_files = list_files(directory)

    if not ppt_files and not pdf_files:
        print("No PPT or PDF files found in the specified directory.")
        return

    print(f"Found {len(ppt_files)} PPT files and {len(pdf_files)} PDF files.")
    choice = input("Do you want to merge (1) PDFs, (2) PPTs, or (3) everything? Enter 1, 2, or 3: ")

    if choice == '1' and pdf_files:
        output_pdf = os.path.join(directory, "merged_pdfs_output.pdf")
        merge_pdfs([os.path.join(directory, pdf) for pdf in pdf_files], output_pdf)
        print(f"All PDFs have been merged into {output_pdf}")
    elif choice == '2' and ppt_files:
        pdf_files = []
        for ppt_file in ppt_files:
            ppt_path = os.path.join(directory, ppt_file)
            pdf_path = os.path.join(directory, ppt_file.replace('.pptx', '.pdf'))
            ppt_to_pdf(ppt_path, pdf_path)
            if os.path.exists(pdf_path):
                pdf_files.append(pdf_path)
            else:
                print(f"Failed to create PDF for {ppt_file}")

        output_pdf = os.path.join(directory, "merged_ppts_output.pdf")
        merge_pdfs(pdf_files, output_pdf)
        if pdf_files:
            print(f"All PPTs have been merged into {output_pdf}")
    elif choice == '3':
        pdf_files = [os.path.join(directory, pdf) for pdf in pdf_files]
        for ppt_file in ppt_files:
            ppt_path = os.path.join(directory, ppt_file)
            pdf_path = os.path.join(directory, ppt_file.replace('.pptx', '.pdf'))
            ppt_to_pdf(ppt_path, pdf_path)
            if os.path.exists(pdf_path):
                pdf_files.append(pdf_path)
            else:
                print(f"Failed to create PDF for {ppt_file}")

        output_pdf = os.path.join(directory, "merged_all_output.pdf")
        merge_pdfs(pdf_files, output_pdf)
        if pdf_files:
            print(f"All files have been merged into {output_pdf}")
    else:
        print("Invalid choice or no files to merge.")

if __name__ == "__main__":
    main()